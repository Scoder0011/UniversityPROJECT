
from flask import Flask, request, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
import tempfile
from pathlib import Path
from datetime import datetime
import json
import time
import subprocess

# Document processing libraries
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from PyPDF2 import PdfMerger, PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from PIL import Image
import io


app = Flask(__name__)
CORS(app)


ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'pptx', 'ppt', 'txt', 'jpg', 'jpeg', 'png', 'gif'}
TEMP_DIR = tempfile.gettempdir()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_type(filename):
    ext = filename.rsplit('.', 1)[1].lower()
    if ext in ['pdf']:
        return 'pdf'
    elif ext in ['docx', 'doc']:
        return 'docx'
    elif ext in ['pptx', 'ppt']:
        return 'pptx'
    elif ext in ['txt']:
        return 'txt'
    elif ext in ['jpg', 'jpeg', 'png', 'gif']:
        return 'image'
    return 'unknown'

# HELPER FUNCTIONS FOR UNIDOC
def create_cover_page(output_path):
    """Create cover page for UniDoc"""
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 36)
    c.drawCentredString(width/2, height - 2*inch, "Course File")
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width/2, height - 3*inch, "Documentation")
    c.setFont("Helvetica", 12)
    c.drawCentredString(width/2, height - 4*inch, f"Generated: {datetime.now().strftime('%B %d, %Y')}")
    c.showPage()
    c.save()

def create_course_info_page(course_data, output_path):
    """Create course information page"""
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 20)
    c.drawString(72, height - 1.5*inch, "Course Information")
    c.setFont("Helvetica-Bold", 12)
    y = height - 2.5*inch
    fields = [
        ('Program & Semester', course_data.get('program', 'N/A')),
        ('Course Code', course_data.get('code', 'N/A')),
        ('Course Name', course_data.get('name', 'N/A')),
        ('Course Coordinator', course_data.get('coordinator', 'N/A')),
        ('Theory Faculty', course_data.get('faculty', 'N/A')),
        ('LTPC', course_data.get('ltpc', 'N/A'))
    ]
    for label, value in fields:
        c.setFont("Helvetica-Bold", 11)
        c.drawString(72, y, f"{label}:")
        c.setFont("Helvetica", 11)
        c.drawString(72, y - 15, str(value))
        y -= 40
    c.showPage()
    c.save()

def create_index_page(file_names, output_path):
    """Create index page"""
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 20)
    c.drawString(72, height - 1.5*inch, "Table of Contents")
    c.setFont("Helvetica", 11)
    y = height - 2.5*inch
    for i, name in enumerate(file_names, 1):
        if y < 100:
            c.showPage()
            y = height - inch
        c.drawString(72, y, f"{i}. {name}")
        y -= 25
    c.showPage()
    c.save()

def create_divider_pdf(title, output_path):
    """Create a one-page divider PDF with title and timestamp."""
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 26)
    c.drawCentredString(width / 2, height - 2 * inch, title)
    c.setFont("Helvetica", 10)
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    c.drawCentredString(width / 2, height - 2 * inch - 20, f"Generated: {timestamp}")
    c.showPage()
    c.save()

# CONVERTERS TO INTERMEDIATE FORMAT
def docx_to_text_with_formatting(docx_path):
    """Extract text with basic formatting info from DOCX"""
    doc = Document(docx_path)
    content = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            content.append({
                'text': para.text,
                'style': para.style.name,
                'alignment': para.alignment
            })
    
    # Extract tables
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        if table_data:
            content.append({'table': table_data})
    
    return content

def pptx_to_text(pptx_path):
    """Extract text from PPTX"""
    prs = Presentation(pptx_path)
    content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = {'slide_num': slide_num, 'shapes': []}
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content['shapes'].append(shape.text)
        content.append(slide_content)
    
    return content

def txt_to_text(txt_path):
    """Read text file"""
    with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def image_to_pdf(image_path, output_pdf):
    """Convert image to PDF properly"""
    img = Image.open(image_path)
    
    # Convert to RGB if necessary
    if img.mode in ('RGBA', 'LA', 'P'):
        background = Image.new('RGB', img.size, (255, 255, 255))
        if img.mode == 'P':
            img = img.convert('RGBA')
        background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
        img = background
    
    # Save as PDF
    img.save(output_pdf, 'PDF', resolution=100.0, quality=95)

def text_to_pdf(text, output_pdf):
    """Convert plain text to PDF with formatting"""
    doc = SimpleDocTemplate(output_pdf, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    for line in text.split('\n'):
        if line.strip():
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 0.1 * inch))
    
    doc.build(story)

def docx_to_pdf(docx_path, output_pdf):
    """Convert DOCX to PDF - try LibreOffice first, fallback to manual"""
    libreoffice_commands = [
        'libreoffice', 'soffice',
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
        'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe'
    ]

    libreoffice_path = None
    for cmd in libreoffice_commands:
        try:
            result = subprocess.run([cmd, '--version'], capture_output=True, timeout=5, text=True)
            if result.returncode == 0:
                libreoffice_path = cmd
                print(f"✅ Found LibreOffice: {cmd}")
                break
        except Exception:
            continue

    # Try LibreOffice if available
    if libreoffice_path:
        try:
            output_dir = os.path.dirname(output_pdf)
            subprocess.run([
                libreoffice_path, '--headless', '--nologo',
                '--nofirststartwizard', '--norestore',
                '--convert-to', 'pdf', '--outdir', output_dir, docx_path
            ], capture_output=True, timeout=60)

            base_name = os.path.splitext(os.path.basename(docx_path))[0]
            temp_pdf = os.path.join(output_dir, f"{base_name}.pdf")

            if os.path.exists(temp_pdf):
                os.rename(temp_pdf, output_pdf)
                print(f"✅ LibreOffice conversion SUCCESS: {os.path.basename(docx_path)}")
                return
        except Exception as e:
            print(f"⚠️ LibreOffice failed, trying manual conversion: {e}")

    # Manual fallback conversion
    print(f"⚠️ Using manual DOCX conversion for: {os.path.basename(docx_path)}")
    try:
        doc = SimpleDocTemplate(output_pdf, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        source_doc = Document(docx_path)
        
        for para in source_doc.paragraphs:
            if para.text.strip():
                story.append(Paragraph(para.text, styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))
        
        # Process tables
        for table in source_doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                ]))
                story.append(Spacer(1, 0.2 * inch))
                story.append(t)
                story.append(Spacer(1, 0.2 * inch))
        
        doc.build(story)
        print(f"✅ Manual conversion successful: {os.path.basename(docx_path)}")
    except Exception as e:
        print(f"❌ Manual DOCX conversion failed: {e}")
        raise

def pptx_to_pdf(pptx_path, output_pdf):
    """Convert PPTX to PDF"""
    doc = SimpleDocTemplate(output_pdf, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontSize=16,
        textColor='#2c3e50',
        spaceAfter=12,
        alignment=1
    )
    
    prs = Presentation(pptx_path)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"Slide {slide_num}", slide_title_style))
        story.append(Spacer(1, 0.2 * inch))
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                story.append(Paragraph(shape.text, styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))
        
        if slide_num < len(prs.slides):
            story.append(PageBreak())
    
    doc.build(story)

# COMBINERS FOR DIFFERENT OUTPUT FORMATS
def combine_to_pdf(files, output_path):
    """Combine all files into a PDF - PRESERVING ORIGINAL PDF FORMATTING"""
    merger = PdfMerger()
    temp_pdfs = []
    
    try:
        for idx, file_info in enumerate(files):
            file_path = file_info['path']
            file_type = file_info['type']
            filename = file_info['name']
            
            if file_type == 'pdf':
                merger.append(file_path)
            else:
                temp_pdf = os.path.join(TEMP_DIR, f"temp_{idx}_{datetime.now().timestamp()}.pdf")
                temp_pdfs.append(temp_pdf)
                
                try:
                    if file_type == 'image':
                        image_to_pdf(file_path, temp_pdf)
                    elif file_type == 'txt':
                        text = txt_to_text(file_path)
                        text_to_pdf(text, temp_pdf)
                    elif file_type == 'docx':
                        docx_to_pdf(file_path, temp_pdf)
                    elif file_type == 'pptx':
                        pptx_to_pdf(file_path, temp_pdf)
                    
                    merger.append(temp_pdf)
                
                except Exception as e:
                    print(f"Error converting {filename}: {e}")
                    error_pdf = os.path.join(TEMP_DIR, f"error_{idx}.pdf")
                    temp_pdfs.append(error_pdf)
                    c = canvas.Canvas(error_pdf, pagesize=letter)
                    c.setFont("Helvetica", 12)
                    c.drawString(100, 750, f"Error processing file: {filename}")
                    c.drawString(100, 730, f"Error: {str(e)}")
                    c.save()
                    merger.append(error_pdf)
        
        merger.write(output_path)
        merger.close()
    
    finally:
        for temp_pdf in temp_pdfs:
            try:
                if os.path.exists(temp_pdf):
                    os.remove(temp_pdf)
            except:
                pass

def combine_to_docx(files, output_path):
    """Combine all files into a DOCX"""
    doc = Document()
    
    title = doc.add_heading('Combined Document', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    for idx, file_info in enumerate(files):
        file_path = file_info['path']
        file_type = file_info['type']
        filename = file_info['name']
        
        heading = doc.add_heading(f'File {idx + 1}: {filename}', level=1)
        heading_para = heading.runs[0]
        heading_para.font.color.rgb = RGBColor(102, 126, 234)
        
        try:
            if file_type == 'pdf':
                reader = PdfReader(file_path)
                for page_num, page in enumerate(reader.pages):
                    doc.add_heading(f'Page {page_num + 1}', level=2)
                    text = page.extract_text()
                    if text.strip():
                        for para in text.split('\n'):
                            if para.strip():
                                doc.add_paragraph(para)
            
            elif file_type == 'docx':
                source_doc = Document(file_path)
                
                for para in source_doc.paragraphs:
                    if para.text.strip():
                        new_para = doc.add_paragraph(para.text)
                        if para.style.name.startswith('Heading'):
                            new_para.style = para.style.name
                
                for table in source_doc.tables:
                    new_table = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                    new_table.style = 'Light Grid Accent 1'
                    
                    for i, row in enumerate(table.rows):
                        for j, cell in enumerate(row.cells):
                            new_table.rows[i].cells[j].text = cell.text
            
            elif file_type == 'pptx':
                content = pptx_to_text(file_path)
                for slide_info in content:
                    doc.add_heading(f"Slide {slide_info['slide_num']}", level=2)
                    for shape_text in slide_info['shapes']:
                        doc.add_paragraph(shape_text)
            
            elif file_type == 'txt':
                text = txt_to_text(file_path)
                for line in text.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line)
            
            elif file_type == 'image':
                try:
                    doc.add_picture(file_path, width=Inches(6))
                except Exception:
                    doc.add_paragraph(f'[Image: {filename}]')
        
        except Exception as e:
            doc.add_paragraph(f'Error processing this file: {str(e)}')
        
        if idx < len(files) - 1:
            doc.add_page_break()
    
    doc.save(output_path)

def combine_to_pptx(files, output_path):
    """Combine all files into a PPTX"""
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Combined Presentation"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    for idx, file_info in enumerate(files):
        file_path = file_info['path']
        file_type = file_info['type']
        filename = file_info['name']
        
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = f'File {idx + 1}: {filename}'
        
        try:
            if file_type == 'pdf':
                reader = PdfReader(file_path)
                for page_num, page in enumerate(reader.pages):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    
                    left = top = PptxInches(0.5)
                    width = PptxInches(9)
                    height = PptxInches(6.5)
                    
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.text = f"Page {page_num + 1}\n\n"
                    
                    text = page.extract_text()
                    if text.strip():
                        tf.text += text[:2000]
            
            elif file_type == 'docx':
                content = docx_to_text_with_formatting(file_path)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                left = top = PptxInches(0.5)
                width = PptxInches(9)
                height = PptxInches(6.5)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                text_content = []
                for item in content:
                    if isinstance(item, dict) and 'text' in item:
                        text_content.append(item['text'])
                
                tf.text = '\n'.join(text_content)[:2000]
            
            elif file_type == 'pptx':
                source_prs = Presentation(file_path)
                for source_slide in source_prs.slides:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    for shape in source_slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            left = PptxInches(0.5)
                            top = PptxInches(0.5)
                            width = PptxInches(9)
                            height = PptxInches(6.5)
                            
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            tf.text = shape.text
            
            elif file_type == 'txt':
                text = txt_to_text(file_path)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                left = top = PptxInches(0.5)
                width = PptxInches(9)
                height = PptxInches(6.5)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.text = text[:2000]
            
            elif file_type == 'image':
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                img = Image.open(file_path)
                img_width, img_height = img.size
                
                max_width = PptxInches(9)
                max_height = PptxInches(6.5)
                
                aspect = img_width / img_height
                if img_width > img_height:
                    width = max_width
                    height = width / aspect
                else:
                    height = max_height
                    width = height * aspect
                
                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 2
                
                slide.shapes.add_picture(file_path, left, top, width, height)
        
        except Exception as e:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = top = PptxInches(1)
            width = PptxInches(8)
            height = PptxInches(5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f'Error processing this file: {str(e)}'
    
    prs.save(output_path)

# ROUTES
@app.route('/combine', methods=['POST'])
def combine_files():
    """Main endpoint to combine files"""
    if 'files' not in request.files:
        return {'error': 'No files uploaded'}, 400
    
    files = request.files.getlist('files')
    output_format = request.form.get('output_format', 'pdf')
    
    if not files or files[0].filename == '':
        return {'error': 'No files selected'}, 400
    
    for file in files:
        if not allowed_file(file.filename):
            return {'error': f'File type not allowed: {file.filename}'}, 400
    
    temp_files = []
    output_path = None
    
    try:
        for file in files:
            filename = secure_filename(file.filename)
            temp_path = os.path.join(TEMP_DIR, f"{datetime.now().timestamp()}_{filename}")
            file.save(temp_path)
            
            temp_files.append({
                'path': temp_path,
                'name': filename,
                'type': get_file_type(filename)
            })
        
        output_filename = f'combined_{datetime.now().strftime("%Y%m%d_%H%M%S")}.{output_format}'
        output_path = os.path.join(TEMP_DIR, output_filename)
        
        if output_format == 'pdf':
            combine_to_pdf(temp_files, output_path)
        elif output_format == 'docx':
            combine_to_docx(temp_files, output_path)
        elif output_format == 'pptx':
            combine_to_pptx(temp_files, output_path)
        else:
            return {'error': 'Invalid output format'}, 400
        
        response = send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype=f'application/{output_format}'
        )
        
        return response
    
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Error: {error_details}")
        return {'error': str(e), 'details': error_details}, 500
    
    finally:
        for file_info in temp_files:
            try:
                if os.path.exists(file_info['path']):
                    os.remove(file_info['path'])
            except:
                pass

@app.route('/combine-checklist', methods=['POST'])
def combine_checklist():
    """Combine files with checklist dividers"""
    if 'checklist_data' not in request.form:
        return {'error': 'Missing checklist_data'}, 400

    try:
        checklist_data = json.loads(request.form['checklist_data'])
    except Exception as e:
        return {'error': 'Invalid checklist_data JSON', 'details': str(e)}, 400

    merger = PdfMerger()
    temp_to_cleanup = []

    try:
        for sec_idx, checklist in enumerate(checklist_data):
            section_name = checklist.get('name') or f"Section {sec_idx+1}"

            divider_fp = os.path.join(TEMP_DIR, f"divider_{int(time.time()*1000)}_{sec_idx}.pdf")
            create_divider_pdf(section_name, divider_fp)
            temp_to_cleanup.append(divider_fp)
            merger.append(divider_fp)

            for file_key in checklist.get('files', []):
                if file_key not in request.files:
                    print(f"Warning: missing file key {file_key}")
                    warn_fp = os.path.join(TEMP_DIR, f"warn_{int(time.time()*1000)}.pdf")
                    c = canvas.Canvas(warn_fp, pagesize=letter)
                    c.setFont("Helvetica", 12)
                    c.drawString(60, 750, f"Missing file for: {file_key}")
                    c.save()
                    temp_to_cleanup.append(warn_fp)
                    merger.append(warn_fp)
                    continue

                fs = request.files[file_key]
                if fs.filename == '':
                    continue

                if not allowed_file(fs.filename):
                    err_fp = os.path.join(TEMP_DIR, f"not_allowed_{int(time.time()*1000)}.pdf")
                    c = canvas.Canvas(err_fp, pagesize=letter)
                    c.setFont("Helvetica", 12)
                    c.drawString(60, 750, f"File type not allowed: {fs.filename}")
                    c.save()
                    temp_to_cleanup.append(err_fp)
                    merger.append(err_fp)
                    continue

                safe_name = secure_filename(fs.filename)
                saved_raw = os.path.join(TEMP_DIR, f"{int(time.time()*1000)}_{safe_name}")
                fs.save(saved_raw)
                temp_to_cleanup.append(saved_raw)

                ftype = get_file_type(safe_name)
                try:
                    if ftype == 'pdf':
                        merger.append(saved_raw)
                    else:
                        converted_pdf = os.path.join(TEMP_DIR, f"converted_{int(time.time()*1000)}.pdf")
                        if ftype == 'image':
                            image_to_pdf(saved_raw, converted_pdf)
                        elif ftype == 'txt':
                            text = txt_to_text(saved_raw)
                            text_to_pdf(text, converted_pdf)
                        elif ftype == 'docx':
                            docx_to_pdf(saved_raw, converted_pdf)
                        elif ftype == 'pptx':
                            pptx_to_pdf(saved_raw, converted_pdf)
                        else:
                            c = canvas.Canvas(converted_pdf, pagesize=letter)
                            c.setFont("Helvetica", 12)
                            c.drawString(60, 750, f"[{safe_name}] - unsupported type, could not convert.")
                            c.save()

                        temp_to_cleanup.append(converted_pdf)
                        merger.append(converted_pdf)

                except Exception as conv_err:
                    error_pdf = os.path.join(TEMP_DIR, f"conv_error_{int(time.time()*1000)}.pdf")
                    c = canvas.Canvas(error_pdf, pagesize=letter)
                    c.setFont("Helvetica", 12)
                    c.drawString(60, 750, f"Error converting file: {safe_name}")
                    c.drawString(60, 730, f"Error: {str(conv_err)}")
                    c.save()
                    temp_to_cleanup.append(error_pdf)
                    merger.append(error_pdf)
                    print(f"Conversion error for {safe_name}:", conv_err)

        output_filename = f'checklist_combined_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pdf'
        output_path = os.path.join(TEMP_DIR, output_filename)
        merger.write(output_path)
        merger.close()

        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype='application/pdf'
        )

    except Exception as e:
        import traceback
        print("ERROR in /combine-checklist:", traceback.format_exc())
        return {'error': str(e)}, 500

    finally:
        for p in temp_to_cleanup:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass

@app.route('/combine-unidoc', methods=['POST'])
def combine_unidoc():
    """Combine files with UniDoc format - cover page, course info, index, then files"""
    files = request.files.getlist("files")
    if not files:
        return {'error': 'No files uploaded'}, 400

    # Collect course metadata from form
    course_data = {
        'program': request.form.get('program', ''),
        'code': request.form.get('code', ''),
        'coordinator': request.form.get('coordinator', ''),
        'name': request.form.get('name', ''),
        'faculty': request.form.get('faculty', ''),
        'ltpc': request.form.get('ltpc', '')
    }

    merger = PdfMerger()
    temp_files = []

    try:
        # Generate the 3 front pages
        cover_fp = os.path.join(TEMP_DIR, f"cover_{int(time.time()*1000)}.pdf")
        info_fp = os.path.join(TEMP_DIR, f"course_info_{int(time.time()*1000)}.pdf")
        index_fp = os.path.join(TEMP_DIR, f"index_{int(time.time()*1000)}.pdf")
        
        create_cover_page(cover_fp)
        create_course_info_page(course_data, info_fp)
        
        # Create index with file names (without extension)
        file_names = [f.filename.rsplit('.', 1)[0] if '.' in f.filename else f.filename for f in files]
        create_index_page(file_names, index_fp)

        # Append front pages
        for fp in [cover_fp, info_fp, index_fp]:
            merger.append(fp)
            temp_files.append(fp)

        # Process uploaded files
        for file in files:
            if file.filename == '':
                continue
                
            filename = secure_filename(file.filename)
            temp_path = os.path.join(TEMP_DIR, f"{datetime.now().timestamp()}_{filename}")
            file.save(temp_path)
            temp_files.append(temp_path)
            
            file_type = get_file_type(filename)
            
            try:
                if file_type == 'pdf':
                    merger.append(temp_path)
                else:
                    converted_pdf = os.path.join(TEMP_DIR, f"converted_{int(time.time()*1000)}.pdf")
                    
                    if file_type == 'image':
                        image_to_pdf(temp_path, converted_pdf)
                    elif file_type == 'txt':
                        text_to_pdf(txt_to_text(temp_path), converted_pdf)
                    elif file_type == 'docx':
                        docx_to_pdf(temp_path, converted_pdf)
                    elif file_type == 'pptx':
                        pptx_to_pdf(temp_path, converted_pdf)
                    else:
                        # Unknown type - create placeholder
                        c = canvas.Canvas(converted_pdf, pagesize=letter)
                        c.setFont("Helvetica", 12)
                        c.drawString(60, 750, f"File: {filename}")
                        c.drawString(60, 730, f"Unsupported file type")
                        c.save()
                    
                    merger.append(converted_pdf)
                    temp_files.append(converted_pdf)
                    
            except Exception as e:
                # Create error page for this file
                error_pdf = os.path.join(TEMP_DIR, f"error_{int(time.time()*1000)}.pdf")
                c = canvas.Canvas(error_pdf, pagesize=letter)
                c.setFont("Helvetica", 12)
                c.drawString(60, 750, f"Error processing file: {filename}")
                c.drawString(60, 730, f"Error: {str(e)}")
                c.save()
                merger.append(error_pdf)
                temp_files.append(error_pdf)
                print(f"Error processing {filename}: {e}")

        # Write final output
        output_filename = f"unidoc_combined_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(TEMP_DIR, output_filename)
        merger.write(output_path)
        merger.close()
        
        return send_file(
            output_path,
            as_attachment=True,
            download_name=output_filename,
            mimetype='application/pdf'
        )

    except Exception as e:
        import traceback
        print("ERROR in /combine-unidoc:", traceback.format_exc())
        return {'error': str(e), 'details': traceback.format_exc()}, 500

    finally:
        for p in temp_files:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except:
                pass

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return {'status': 'healthy', 'timestamp': datetime.now().isoformat()}

@app.route('/', methods=['GET'])
def index():
    """Root endpoint with API info"""
    return {
        'service': 'Universal File Combiner API',
        'version': '2.0.0',
        'endpoints': {
            '/combine': 'POST - Combine multiple files (preserves PDF formatting)',
            '/combine-checklist': 'POST - Combine files with divider pages',
            '/combine-unidoc': 'POST - Create UniDoc with cover, info, and index pages',
            '/health': 'GET - Health check'
        },
        'supported_formats': list(ALLOWED_EXTENSIONS),
        'output_formats': ['pdf', 'docx', 'pptx'],
        'features': [
            'Preserves original PDF formatting (tables, layouts, fonts)',
            'Direct PDF merging without text extraction',
            'Smart conversion for other formats',
            'Checklist mode with section dividers',
            'UniDoc builder for course documentation'
        ]
    }

if __name__ == '__main__':
    print("=" * 60)
    print("Universal File Combiner Backend v2.0")
    print("=" * 60)
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=False, host='0.0.0.0', port=port)
